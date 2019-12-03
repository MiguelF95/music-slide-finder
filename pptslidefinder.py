import glob
import string
from pptx import Presentation

songDict = {}
inputStr = []

print("\n\nThis script will take in PowerPoint Slide titles and return them with their corresponding slide numbers.\n")


#Obtain Slide titles from user input and add to a dictionary with value -1. Input requests stop when user inputs 1
while '1' not in inputStr:
    #In the case that user input all titles at the same time via copy and paste, split into separate strings and store into a vector.
    #Otherwise, each song is put into the vector one at a time.
    inputStr = (input("Input a slide title and press enter to continue adding more titles to the list.\nWhen finished, type '1' and press enter\n\n")).split("\n")
    
    #Create keys using input string, removing punctuation and setting all to lower case.
    for mystr in inputStr:
        key = mystr.translate(mystr.maketrans('', '', string.punctuation)).lower().replace(" ", "")
        songDict[key] = [mystr,-1]
    print("\n\n")


#Remove 1 from songlist
del songDict['1']

#Get Powerpoint File from current directory (same directory as this script)
pptname = glob.glob('*.pptx')
prs = Presentation(pptname[0])


#if a slide's title matches one in the songlist, update the value to the corresponding slide number
for slide in prs.slides:
    keyCheck = slide.shapes.title.text.translate(str.maketrans('', '', string.punctuation)).lower().replace(" ", "")
    if keyCheck in songDict:
        songDict[keyCheck] = [slide.shapes.title.text,
                              prs.slides.index(slide) + 1]


#Print slide numbers with their corresponding titles
#Titles that were not found remain -1 in dictionary and  are printed afterwards as Not Found
print("Slide Numbers from " + pptname[0])
print("=============================================================")

for i in songDict:
    if songDict[i][1] != -1:
         print(str(songDict[i]))

for i in songDict:
    if songDict[i][1] == -1:
         print('Not Found: ' + songDict[i][0])

print("=============================================================")
